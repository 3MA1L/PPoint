import os
import datetime
import collections
import collections.abc
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from enum import Enum
from pptx.util import Inches
from pptx.chart.data import ChartData
from comtypes.client import Constants, CreateObject

def update_placeholders(prs, placeholders):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for placeholder, value in placeholders.items():
                            run.text = run.text.replace(placeholder, value)

def create_powerpoint_slide(prs, slide_layout, placeholders):
    slide = prs.slides.add_slide(slide_layout)
    for placeholder, text in placeholders.items():
        for shape in slide.shapes:
            if shape.has_text_frame and placeholder in shape.text_frame.text:
                shape.text_frame.text = text

#def create_bar_chart(prs, incidents_last_call, medium_incidents_last_call, low_incidents_last_call,
                     # incidents_current_call, medium_incidents_current_call, low_incidents_current_call):
    # slide_layout = prs.slide_layouts[5]  # Use index 5 for the title and content slide layout with a chart
    # slide = prs.slides.add_slide(slide_layout)

    # x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
    # chart = slide.shapes.add_chart(
        # XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy
    # ).chart

    #Add chart data
    # chart_data = CategoryChartData()
    # chart_data.categories = ['Incidents', 'Medium Incidents', 'Low Incidents']
    # chart_data.add_series('Last Call', (incidents_last_call, medium_incidents_last_call, low_incidents_last_call))
    # chart_data.add_series('Current Call', (incidents_current_call, medium_incidents_current_call, low_incidents_current_call))

    # chart.replace_data(chart_data)
####
def save_powerpoint_as_pdf(prs, output_filename):
    pdf_export_filename = output_filename.replace(".pptx", ".pdf")
    powerpoint = CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1
    ppt = powerpoint.Presentations.Open(output_filename)
    ppt.SaveAs(pdf_export_filename, Constants.ppSaveAsPDF)
    ppt.Close()
    powerpoint.Quit()
    
def set_named_placeholder(slide, shape_name, new_value):
    #Find the shape byt its name in the slide shapes
    for shape in slide.shapes:
        if shape.name == shape_name:
            for paragrap in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.text = run.text.replace(placeholder_name, new_value)
            
            break

def main():
    # Load the PowerPoint template
    script_dir = os.path.dirname(os.path.abspath(__file__))
    template_filename = os.path.join(script_dir, "SecuritySlides.pptx")
    
    prs = Presentation(template_filename)
    
    slide_master = prs.slide_master

    # Slide 1: Ask for customer name and update the template

    customer_name = input("Enter Customer Name: ")
    placeholders_slide1 = {
        'CUSTOMER_NAME': customer_name
    }
    set_named_placeholder(slide_master, "CUSTOMER_NAME", customer_name)
    update_placeholders(prs, placeholders_slide1)
###########################################################################################################################################
    # Slide 2: Ask for follow-up items and update the template
    follow_up_items_list = []
    
    while True:
        follow_up_items = input("Enter a follow up item: ")
        
        #if the user enters 'done', exit the loop
        if follow_up_items.lower() == 'done':
            break
        
            
        follow_up_items_list.append(follow_up_items)
        
        follow_up = '\n'.join(follow_up_items_list)
    
    placeholders_slide2 = {
        'FOLLOW_UP_ITEMS': follow_up
    }
    set_named_placeholder(slide_master, 'FOLLOW_UP_ITEMS', follow_up)
    update_placeholders(prs, placeholders_slide2)
    ########################################################################################################################################
    # Slide 3: Ask for incident details and update the template
    incidents_last_call = int(input("Enter the number of Incidents for Last Call: "))
    medium_incidents_last_call = int(input("Enter the number of Medium Incidents for Last Call: "))
    low_incidents_last_call = int(input("Enter the number of Low Incidents for Last Call: "))

    incidents_current_call = int(input("Enter the number of Incidents for Current Call: "))
    medium_incidents_current_call = int(input("Enter the number of Medium Incidents for Current Call: "))
    low_incidents_current_call = int(input("Enter the number of Low Incidents for Current Call: "))

    placeholders_slide3 = {
        'INCIDENTS_LAST_CALL': str(incidents_last_call),
        'MEDIUM_INCIDENTS_LAST_CALL': str(medium_incidents_last_call),
        'LOW_INCIDENTS_LAST_CALL': str(low_incidents_last_call),
        'INCIDENTS_CURRENT_CALL': str(incidents_current_call),
        'MEDIUM_INCIDENTS_CURRENT_CALL': str(medium_incidents_current_call),
        'LOW_INCIDENTS_CURRENT_CALL': str(low_incidents_current_call)
    }
    update_placeholders(prs, placeholders_slide3)
    
 ###############################################################################################################################################################
 #chart for slide3
    slide_layout = prs.slide_layouts[5] 
    slide = prs.slides.add_slide(slide_layout)
    
    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy,
    ).chart
    
 #ask the user for chart data
    num_categories = int(input("Enter the number of categories for the chart: "))
    categories = [input(f"Enter category {i+1}: ") for i in range(num_categories)]
    
    num_series = int(input("Enter the number of series for the chart: "))
    chart_data = []
    
    for i in range(num_series):
        series_name = input(f"Enter series {i+1} name: ")
        data_values = [int(input(f"Enter data value for '{series_name}' in category '{categories[j]}: ")) for j in range(num_categories)]
        
        chart_data.add_series(series_name, data,values)
        
    chart_data.categories = categories   
    chart.chart_data.add_series(chart_data)
    
    
        
#################################################################################################################################################
    # Save the PowerPoint presentation with the desired filename format
    current_date = datetime.datetime.now().strftime("%Y-%m-%d")
    output_folder = os.path.join(os.path.expanduser("~"), "Desktop", customer_name)
    os.makedirs(output_folder, exist_ok=True)
    output_filename = os.path.join(output_folder, f"SecurityCallSlides_{customer_name}_{current_date}.pptx")
    prs.save(output_filename)
    print(f"PowerPoint presentation '{output_filename}' created successfully.")

    # Export the PowerPoint presentation as PDF
    pdf_output_folder = os.path.join(output_folder, "PDFs")
    os.makedirs(pdf_output_folder, exist_ok=True)
    save_powerpoint_as_pdf(prs, output_filename)
    pdf_export_filename = output_filename.replace(".pptx", ".pdf")
    pdf_output_filename = os.path.join(pdf_output_folder, os.path.basename(pdf_export_filename))
    os.rename(pdf_export_filename, pdf_output_filename)
    print(f"PDF version '{pdf_output_filename}' created successfully.")

if __name__ == "__main__":
    main()